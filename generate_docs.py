import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

def generate_pdf(output_path):
    pdf = FPDF()
    pdf.add_page()
    
    # Title
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, "SmartHire Project Report", ln=True, align="C")
    pdf.ln(10)
    
    # Content
    pdf.set_font("Arial", '', 12)
    content = [
        "1. Overview",
        "SmartHire is an AI-Powered Resume-to-Job Matching and Career Guidance Engine. It uses classical Machine Learning models to match uploaded resumes against a corpus of jobs, predict job categories, and identify skill gaps.",
        "",
        "2. Architecture & Tech Stack",
        "- Language: Python 3.10+",
        "- Web Framework: Streamlit",
        "- Data Processing: Pandas, NumPy, Scikit-Learn",
        "- NLP & Machine Learning: TF-IDF, Logistic Regression, Cosine Similarity, KMeans",
        "",
        "3. Core Components",
        "- Resume Classifier: Supervised model using TF-IDF and Logistic Regression to predict the best job category for a given resume.",
        "- Job Recommender: Unsupervised matching engine using TF-IDF and Cosine Similarity to find top job matches based on skill overlap.",
        "- Skill-Gap Analysis: Compares extracted skills from the resume with required skills from job postings.",
        "",
        "4. Workflow",
        "- Data loading and preprocessing from Kaggle datasets.",
        "- Model training via Jupyter notebooks.",
        "- Inference served through an interactive Streamlit frontend."
    ]
    
    for line in content:
        pdf.multi_cell(0, 8, line)
        
    pdf.output(output_path)
    print(f"PDF successfully generated at {output_path}")

def generate_ppt(output_path):
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SmartHire"
    subtitle.text = "AI-Powered Resume-to-Job Matching & Career Guidance Engine"
    
    # Architecture Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "System Architecture"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Core Components:"
    
    p = tf.add_paragraph()
    p.text = "1. Resume Parsing & Skill Extraction"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Category Classification (Logistic Regression)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Job Recommendation (Cosine Similarity)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. Streamlit Web Portal"
    p.level = 1

    # Save
    prs.save(output_path)
    print(f"PPTX successfully generated at {output_path}")

if __name__ == "__main__":
    os.makedirs("reports", exist_ok=True)
    pdf_path = os.path.join("reports", "SmartHire_Project_Report.pdf")
    ppt_path = os.path.join("reports", "SmartHire_Project_Presentation.pptx")
    
    try:
        generate_pdf(pdf_path)
    except Exception as e:
        print(f"Failed to generate PDF: {e}")
        
    try:
        generate_ppt(ppt_path)
    except Exception as e:
        print(f"Failed to generate PPT: {e}")
